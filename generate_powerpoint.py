"""
Script de génération de la présentation PowerPoint (PPTX) pour le projet Molly Market.
Intègre le design complet, l'architecture PostgreSQL, la modélisation Merise,
et les captures d'écran des tests de chaque rôle.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    
    # Palette de couleurs Molly Market
    NAVY = RGBColor(26, 35, 126)     # #1A237E
    ORANGE = RGBColor(251, 140, 0)   # #FB8C00
    DARK = RGBColor(33, 33, 33)      # #212121
    LIGHT_BG = RGBColor(245, 247, 250)
    WHITE = RGBColor(255, 255, 255)
    GREEN = RGBColor(46, 125, 50)
    GRAY = RGBColor(117, 117, 117)
    
    screenshots_dir = r"C:\Users\ASUS\.gemini\antigravity-ide\brain\0ce5b187-e2fa-4a7a-98da-c40318458593"
    
    def add_header(slide, title_text, category_text="MOLLY MARKET — SOUTENANCE DE PROJET"):
        # Header banner
        header_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.1))
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = NAVY
        header_shape.line.color.rgb = NAVY
        
        # Category subtitle
        tx_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.12), Inches(11), Inches(0.3))
        tf = tx_box.text_frame
        p = tf.paragraphs[0]
        p.text = category_text.upper()
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = ORANGE
        
        # Slide Title
        tx_box2 = slide.shapes.add_textbox(Inches(0.8), Inches(0.38), Inches(11), Inches(0.6))
        tf2 = tx_box2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = title_text
        p2.font.size = Pt(22)
        p2.font.bold = True
        p2.font.color.rgb = WHITE
        
        # Footer
        footer_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.1), Inches(13.333), Inches(0.4))
        footer_shape.fill.solid()
        footer_shape.fill.fore_color.rgb = RGBColor(230, 235, 245)
        footer_shape.line.fill.background()
        
        ft_box = slide.shapes.add_textbox(Inches(0.8), Inches(7.12), Inches(11.5), Inches(0.3))
        ft_p = ft_box.text_frame.paragraphs[0]
        ft_p.text = "Molly Market — Plateforme Web & Backend PostgreSQL | Évaluation Finale SQL"
        ft_p.font.size = Pt(10)
        ft_p.font.color.rgb = GRAY

    # -------------------------------------------------------------
    # SLIDE 1 : TITRE
    # -------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = NAVY
    bg1.line.color.rgb = NAVY
    
    # Accent badge
    badge = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(1.2), Inches(4.5), Inches(0.5))
    badge.fill.solid()
    badge.fill.fore_color.rgb = ORANGE
    badge.line.color.rgb = ORANGE
    bp = badge.text_frame.paragraphs[0]
    bp.text = "PROJET D'ÉVALUATION FINALE SQL"
    bp.font.size = Pt(12)
    bp.font.bold = True
    bp.font.color.rgb = WHITE
    bp.alignment = PP_ALIGN.CENTER
    
    # Main Title
    tbox = slide1.shapes.add_textbox(Inches(1.2), Inches(1.9), Inches(11), Inches(2.2))
    tf = tbox.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = "MOLLY MARKET"
    p1.font.size = Pt(44)
    p1.font.bold = True
    p1.font.color.rgb = WHITE
    
    p2 = tf.add_paragraph()
    p2.text = "Plateforme Web de Gestion de Supermarché & Cerveau Métier PostgreSQL"
    p2.font.size = Pt(22)
    p2.font.color.rgb = RGBColor(200, 220, 255)
    
    # Details card
    card = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(4.5), Inches(10.9), Inches(2.2))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(38, 50, 160)
    card.line.color.rgb = RGBColor(60, 80, 200)
    ctf = card.text_frame
    ctf.word_wrap = True
    
    cp1 = ctf.paragraphs[0]
    cp1.text = "Architecture : 100% Logique Métier & Règles de Gestion dans PostgreSQL (PL/pgSQL)"
    cp1.font.size = Pt(14)
    cp1.font.bold = True
    cp1.font.color.rgb = ORANGE
    
    cp2 = ctf.add_paragraph()
    cp2.text = "• SGBD : PostgreSQL 16 (Procédures Stockées ACID, Triggers, Vues, Rôles RBAC)"
    cp2.font.size = Pt(13)
    cp2.font.color.rgb = WHITE
    
    cp3 = ctf.add_paragraph()
    cp3.text = "• Modélisation : Méthode MERISE (Dictionnaire, SAT, MCD, MLD en 3FN)"
    cp3.font.size = Pt(13)
    cp3.font.color.rgb = WHITE
    
    cp4 = ctf.add_paragraph()
    cp4.text = "• Interface : React / Vite connecté via API technique Express (sans logique métier)"
    cp4.font.size = Pt(13)
    cp4.font.color.rgb = WHITE

    # -------------------------------------------------------------
    # SLIDE 2 : CONTEXTE & OBJECTIFS DU PROJET
    # -------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    add_header(slide2, "1. Contexte, Enjeux & Contraintes Techniques")
    
    # Left Card : Besoins Métier
    c_left = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.4), Inches(5.6), Inches(5.4))
    c_left.fill.solid()
    c_left.fill.fore_color.rgb = WHITE
    c_left.line.color.rgb = RGBColor(220, 225, 235)
    
    ltf = c_left.text_frame
    ltf.word_wrap = True
    lp0 = ltf.paragraphs[0]
    lp0.text = "🎯 Besoins Métier du Supermarché"
    lp0.font.size = Pt(16)
    lp0.font.bold = True
    lp0.font.color.rgb = NAVY
    
    items_left = [
        "Gestion globale du catalogue (126 produits réels, 16 rayons/catégories).",
        "Encaissement rapide en caisse (POS avec scanner code-barre, Wave, Espèces).",
        "Suivi des stocks en temps réel avec alertes de rupture automatiques.",
        "Gestion du réapprovisionnement & commandes fournisseurs (Achats).",
        "Tableaux de bord d'aide à la décision pour la Direction.",
        "Traçabilité totale des mouvements de stock et des clôtures de caisse."
    ]
    for item in items_left:
        p = ltf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(12)
        p.font.color.rgb = DARK
        p.space_after = Pt(8)
        
    # Right Card : Contraintes d'architecture
    c_right = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.4), Inches(5.7), Inches(5.4))
    c_right.fill.solid()
    c_right.fill.fore_color.rgb = RGBColor(240, 245, 255)
    c_right.line.color.rgb = RGBColor(180, 200, 240)
    
    rtf = c_right.text_frame
    rtf.word_wrap = True
    rp0 = rtf.paragraphs[0]
    rp0.text = "⚙️ Règle d'Or de l'Architecture"
    rp0.font.size = Pt(16)
    rp0.font.bold = True
    rp0.font.color.rgb = ORANGE
    
    items_right = [
        "100% de la logique métier réside dans PostgreSQL : calculs, validations, stocks, décrets.",
        "L'application Web est un client léger : uniquement formulaires, affichage et appel des procédures SQL.",
        "Intégrité garantie au niveau SGBD : même hors du Web, la base est inviolable.",
        "Transactions ACID : aucune vente ou achat à moitié validé.",
        "Sécurité RBAC : 4 profils utilisateurs étanches avec privilèges spécifiques."
    ]
    for item in items_right:
        p = rtf.add_paragraph()
        p.text = "✔ " + item
        p.font.size = Pt(12)
        p.font.color.rgb = NAVY
        p.space_after = Pt(10)

    # -------------------------------------------------------------
    # SLIDE 3 : CONCEPTION MERISE (MCD & MLD)
    # -------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    add_header(slide3, "2. Modélisation de Données (MCD & MLD Normalisé)")
    
    # 4 columns/cards representing the 4 core pillars
    cols = [
        ("CLIENT & VENTE", ["CLIENT (id_client, nom, prenom...)", "VENTE (id_vente, numero_ticket...)", "LIGNE_VENTE (id_vente, id_produit)", "PAIEMENT (id_paiement, montant, mode)"], NAVY),
        ("PRODUIT & STOCK", ["CATEGORIE (id_categorie, code, nom)", "PRODUIT (id_produit, code_barre, prix)", "MOUVEMENT_STOCK (id_mouv, sens)", "HISTORIQUE_PRIX (ancien, nouveau)"], ORANGE),
        ("FOURNISSEUR & ACHAT", ["FOURNISSEUR (id_fournisseur, nom)", "ACHAT (id_achat, numero_achat, statut)", "LIGNE_ACHAT (id_achat, id_produit)", "FACTURE_FOURNISSEUR (ref, statut)"], GREEN),
        ("SÉCURITÉ & CAISSE", ["UTILISATEUR (id, role, hash...)", "POINTS_CAISSE (session, ecart)", "BILLETAGE_CAISSE (theorique, compte)", "JOURNAL_SUPPRESSIONS (audit jsonb)"], DARK)
    ]
    
    for i, (col_title, col_items, col_color) in enumerate(cols):
        x = Inches(0.8 + i * 2.95)
        card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.4), Inches(2.8), Inches(5.4))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = col_color
        card.line.width = Pt(2)
        
        ctf = card.text_frame
        ctf.word_wrap = True
        cp = ctf.paragraphs[0]
        cp.text = col_title
        cp.font.size = Pt(13)
        cp.font.bold = True
        cp.font.color.rgb = col_color
        
        for it in col_items:
            p = ctf.add_paragraph()
            p.text = "• " + it
            p.font.size = Pt(10)
            p.font.color.rgb = DARK
            p.space_after = Pt(6)
            
    # -------------------------------------------------------------
    # SLIDE 4 : LE COEUR POSTGRESQL (PROCÉDURES, TRIGGERS, VUES)
    # -------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    add_header(slide4, "3. Le Cerveau Métier dans PostgreSQL (PL/pgSQL)")
    
    pillars = [
        ("Procédures Stockées (ACID)", [
            "effectuer_vente(...) : création ticket, décrémentation stock, insertion lignes & paiement.",
            "effectuer_achat(...) : création commande fournisseur.",
            "reception_stock(...) : validation livraison & incrémentation stock.",
            "payer_facture_fournisseur(...) : validation décaissement directeur."
        ], Inches(1.4)),
        ("Déclencheurs (Triggers)", [
            "trg_fn_calcul_montant_commande : calcule et met à jour automatiquement le total de vente.",
            "trg_fn_maj_stock_vente : crée le mouvement de stock et ajuste le stock actuel en temps réel.",
            "trg_fn_empecher_stock_negatif : bloque toute tentative de vente si stock insuffisant.",
            "trg_fn_historique_prix & journal_suppression : trace les modifications et suppressions."
        ], Inches(3.2)),
        ("Vues & Fonctions Analytiques", [
            "vue_statistiques & vue_stock : métriques dashboard, alertes ruptures.",
            "vue_top_produits & vue_top_clients : classements basés sur chiffre_affaires.",
            "chiffre_affaires(periode) : fonction PL/pgSQL pour calculs journaliers/mensuels."
        ], Inches(5.1))
    ]
    
    for title, items, y_pos in pillars:
        c = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y_pos, Inches(11.7), Inches(1.6))
        c.fill.solid()
        c.fill.fore_color.rgb = WHITE
        c.line.color.rgb = RGBColor(200, 215, 235)
        
        tf = c.text_frame
        tf.word_wrap = True
        p0 = tf.paragraphs[0]
        p0.text = title
        p0.font.size = Pt(13)
        p0.font.bold = True
        p0.font.color.rgb = NAVY
        
        for item in items:
            p = tf.add_paragraph()
            p.text = "• " + item
            p.font.size = Pt(10.5)
            p.font.color.rgb = DARK

    # -------------------------------------------------------------
    # SLIDE 5 : SÉCURITÉ RBAC & 4 ACTEURS
    # -------------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    add_header(slide5, "4. Contrôle d'Accès & Matrice des Privilèges RBAC")
    
    roles = [
        ("👑 Administrateur", "admin@mollymarket.ci", "Accès complet", "Gestion des employés, configuration système, accès à tous les modules et journaux d'audit.", NAVY),
        ("💼 Directeur", "directeur@mollymarket.ci", "Validation & Pilotage", "Validation des commandes fournisseurs, approbation des clôtures caisse, analyse du CA et bilans financiers.", ORANGE),
        ("📦 Magasinier", "magasinier@mollymarket.ci", "Gestion des Stocks", "Réception des marchandises fournisseurs, ajustements d'inventaire, alertes de réapprovisionnement.", GREEN),
        ("🛒 Vendeur / Caissier", "vendeur@mollymarket.ci", "Point de Vente (POS)", "Scan code-barre, enregistrement des tickets, encaissement multi-moyens (Wave, Espèces), remise client.", DARK)
    ]
    
    for i, (r_title, r_email, r_perm, r_desc, r_color) in enumerate(roles):
        y = Inches(1.4 + i * 1.35)
        card = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(11.7), Inches(1.2))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = r_color
        card.line.width = Pt(1.5)
        
        tf = card.text_frame
        tf.word_wrap = True
        p0 = tf.paragraphs[0]
        p0.text = f"{r_title}  |  Identifiant : {r_email}  |  Niveau : {r_perm}"
        p0.font.size = Pt(12)
        p0.font.bold = True
        p0.font.color.rgb = r_color
        
        p1 = tf.add_paragraph()
        p1.text = r_desc
        p1.font.size = Pt(11)
        p1.font.color.rgb = DARK

    # Helper function to add test screenshot slides
    def add_screenshot_slide(title, subtitle, img_filename, bullet_points):
        slide = prs.slides.add_slide(blank_layout)
        add_header(slide, title, subtitle)
        
        # Left side : Screenshot
        img_path = os.path.join(screenshots_dir, img_filename)
        if os.path.exists(img_path):
            slide.shapes.add_picture(img_path, Inches(0.8), Inches(1.35), Inches(7.5), Inches(5.4))
        else:
            ph = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.35), Inches(7.5), Inches(5.4))
            ph.fill.solid()
            ph.fill.fore_color.rgb = RGBColor(230, 230, 230)
            ph.text_frame.paragraphs[0].text = f"Capture : {img_filename}"
            
        # Right side : Explanatory card
        rcard = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.5), Inches(1.35), Inches(4.0), Inches(5.4))
        rcard.fill.solid()
        rcard.fill.fore_color.rgb = WHITE
        rcard.line.color.rgb = RGBColor(220, 225, 235)
        
        rtf = rcard.text_frame
        rtf.word_wrap = True
        rp0 = rtf.paragraphs[0]
        rp0.text = "Points Clés Validés"
        rp0.font.size = Pt(14)
        rp0.font.bold = True
        rp0.font.color.rgb = NAVY
        
        for bp in bullet_points:
            p = rtf.add_paragraph()
            p.text = "✔ " + bp
            p.font.size = Pt(11)
            p.font.color.rgb = DARK
            p.space_after = Pt(8)

    # -------------------------------------------------------------
    # SLIDE 6 : TEST ADMINISTRATEUR
    # -------------------------------------------------------------
    add_screenshot_slide(
        "5. Test Acteur : Administrateur (Catalogue & Audit)",
        "VALIDATION EN DIRECT — RÔLE ADMINISTRATEUR",
        "admin_dashboard_1788564187203.png",
        [
            "Connexion réussie avec `admin@mollymarket.ci`.",
            "Chargement instantané des statistiques du tableau de bord.",
            "126 produits réels actifs répartis sur 16 catégories/rayons.",
            "Gestion des 4 comptes employés et de leurs habilitations RBAC.",
            "Consultation du journal d'audit des mouvements de stock."
        ]
    )

    # -------------------------------------------------------------
    # SLIDE 7 : TEST DIRECTEUR
    # -------------------------------------------------------------
    add_screenshot_slide(
        "6. Test Acteur : Directeur (Achats & Décisionnel)",
        "VALIDATION EN DIRECT — RÔLE DIRECTEUR",
        "directeur_achats_1788564392421.png",
        [
            "Connexion avec `directeur@mollymarket.ci`.",
            "Workflow de validation des bons de commande fournisseurs.",
            "Approbation des dépenses et décaissements de trésorerie.",
            "Visualisation des courbes d'évolution du chiffre d'affaires (FCFA).",
            "Export des rapports statistiques aux formats PDF et Excel."
        ]
    )

    # -------------------------------------------------------------
    # SLIDE 8 : TEST MAGASINIER
    # -------------------------------------------------------------
    add_screenshot_slide(
        "7. Test Acteur : Magasinier (Stocks & Réceptions)",
        "VALIDATION EN DIRECT — RÔLE MAGASINIER",
        "magasinier_stocks_1788564532097.png",
        [
            "Connexion avec `magasinier@mollymarket.ci`.",
            "Surveillance des niveaux de stocks et alertes de seuil critique.",
            "Réception physique des livraisons fournisseurs (`reception_stock`).",
            "Incrémentation automatique des quantités en stock via procédure SQL.",
            "Génération immédiate d'une trace d'audit d'entrée de stock."
        ]
    )

    # -------------------------------------------------------------
    # SLIDE 9 : TEST VENDEUR / CAISSIER
    # -------------------------------------------------------------
    add_screenshot_slide(
        "8. Test Acteur : Vendeur (POS & Encaissement)",
        "VALIDATION EN DIRECT — RÔLE VENDEUR / CAISSIER",
        "vendeur_receipt_modal_1788564924761.png",
        [
            "Connexion avec `vendeur@mollymarket.ci`.",
            "Ajout d'articles au panier et calcul automatique de la TVA (18%).",
            "Encaissement en espèces avec 10 000 FCFA remis (Monnaie : 1 400 FCFA).",
            "Exécution de la procédure SQL `effectuer_vente(...)`.",
            "Génération et émission du ticket de caisse officiel `TK-20260905-001`."
        ]
    )

    # -------------------------------------------------------------
    # SLIDE 10 : CONCLUSION & PERSPECTIVES
    # -------------------------------------------------------------
    slide10 = prs.slides.add_slide(blank_layout)
    add_header(slide10, "9. Conclusion & Points Forts du Projet")
    
    # 3 Summary boxes
    strengths = [
        ("🏆 Respect Total du Cahier des Charges", [
            "100% de la logique métier implémentée dans PostgreSQL (PL/pgSQL).",
            "Client Web ultra-léger sans duplication de règles de calcul.",
            "Base de données normalisée (3FN) avec 11 tables interconnectées."
        ], NAVY),
        ("🔒 Sécurité & Robustesse Entreprise", [
            "Transactions ACID garantissant zéro incohérence de stock ou de caisse.",
            "Contrôle d'accès basé sur les rôles (RBAC à 4 profils).",
            "Journalisation d'audit automatique des modifications et suppressions."
        ], GREEN),
        ("🚀 Évolutivité & Performance", [
            "Calculs exécutés au plus près de la donnée (mémoire PostgreSQL).",
            "126 produits réels avec photos et codes-barres pré-chargés.",
            "Prêt pour le déploiement en réseau multi-caisses."
        ], ORANGE)
    ]
    
    for i, (st_title, st_items, st_color) in enumerate(strengths):
        y = Inches(1.4 + i * 1.8)
        box = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(11.7), Inches(1.6))
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = st_color
        box.line.width = Pt(2)
        
        btf = box.text_frame
        btf.word_wrap = True
        bp0 = btf.paragraphs[0]
        bp0.text = st_title
        bp0.font.size = Pt(13)
        bp0.font.bold = True
        bp0.font.color.rgb = st_color
        
        for it in st_items:
            p = btf.add_paragraph()
            p.text = "✔ " + it
            p.font.size = Pt(10.5)
            p.font.color.rgb = DARK

    # Save presentation
    output_path = r"c:\Users\ASUS\mollymarket\MollyMarket_Presentation_Soutenance.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")

if __name__ == "__main__":
    create_presentation()
